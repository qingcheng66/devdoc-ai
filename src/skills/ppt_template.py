"""PPT Template Engine — fills structured content into .pptx slides.

Two modes:
1. External template: .pptx with {{placeholder}} text → detect type → clone → fill
2. Built-in: generates professional slides via python-pptx shapes directly
"""

import copy
import io
import re
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from src.utils.logger import logger

PLACEHOLDER_RE = re.compile(r"\{\{(\w+)\}\}")

# ── Color Scheme ──────────────────────────────────────────────────
C_DARK_BG   = RGBColor(0x1a, 0x36, 0x5d)   # deep blue
C_LIGHT_BG  = RGBColor(0xf5, 0xf5, 0xf5)   # light gray
C_WHITE     = RGBColor(0xff, 0xff, 0xff)
C_DARK_TEXT = RGBColor(0x1a, 0x20, 0x2c)
C_ACCENT    = RGBColor(0x3b, 0x82, 0xf6)   # blue
C_ACCENT2   = RGBColor(0xf5, 0x9e, 0x0b)   # amber
C_LIGHT_BLUE = RGBColor(0xbf, 0xdb, 0xfe)
C_GRAY_TEXT = RGBColor(0x9c, 0xa3, 0xaf)

FONT_NAME = "Microsoft YaHei"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_font(run, size_pt: int, bold: bool = False, color=None, name: str = FONT_NAME):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = name
    if color:
        run.font.color.rgb = color


def _add_textbox(slide, left, top, width, height) -> "TextFrame":
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# ── Built-in Template ─────────────────────────────────────────────

class BuiltinTemplate:
    """Generate a complete .pptx with professional built-in styling."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    def add_cover(self, title: str, subtitle: str = ""):
        slide = _blank_slide(self.prs)
        _set_bg(slide, C_DARK_BG)

        # Accent line at top
        slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_W, Inches(0.06)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = C_ACCENT
        slide.shapes[-1].line.fill.background()

        tf = _add_textbox(slide, 1.5, 2.2, 10.333, 1.8)
        p = tf.paragraphs[0]
        p.text = title
        _set_font(p.runs[0] if p.runs else p.add_run(), 52, True, C_WHITE)
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            tf2 = _add_textbox(slide, 1.5, 4.2, 10.333, 1.2)
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            _set_font(p2.runs[0] if p2.runs else p2.add_run(), 26, False, C_LIGHT_BLUE)
            p2.alignment = PP_ALIGN.CENTER

        # Bottom accent bar
        slide.shapes.add_shape(
            1, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = C_ACCENT
        slide.shapes[-1].line.fill.background()

    def add_toc(self, title: str, items: list[str]):
        slide = _blank_slide(self.prs)
        _set_bg(slide, C_WHITE)

        tf = _add_textbox(slide, 1.0, 0.6, 11.333, 1.0)
        p = tf.paragraphs[0]
        p.text = title
        _set_font(p.runs[0] if p.runs else p.add_run(), 44, True, C_DARK_BG)

        # Underline
        slide.shapes.add_shape(
            1, Inches(1.0), Inches(1.55), Inches(2.5), Inches(0.04)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = C_ACCENT
        slide.shapes[-1].line.fill.background()

        tf2 = _add_textbox(slide, 1.5, 2.2, 10.333, 4.8)
        for i, item in enumerate(items):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"  {item}"
            _set_font(p.runs[0] if p.runs else p.add_run(), 28, False, C_DARK_TEXT)
            p.space_after = Pt(18)
            # Number prefix
            p.text = f"  {i + 1:02d}    {item}"
            p.runs[0].font.color.rgb = C_ACCENT if p.runs else C_DARK_TEXT

    def add_content(self, title: str, body: str, slide_num: int = 0, total: int = 0):
        slide = _blank_slide(self.prs)
        _set_bg(slide, C_WHITE)

        # Header bar
        slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_W, Inches(1.2)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = C_DARK_BG
        slide.shapes[-1].line.fill.background()

        tf = _add_textbox(slide, 1.0, 0.15, 10.5, 1.0)
        p = tf.paragraphs[0]
        p.text = title
        _set_font(p.runs[0] if p.runs else p.add_run(), 40, True, C_WHITE)

        # Page number
        if total:
            tf_pn = _add_textbox(slide, 11.5, 0.45, 1.5, 0.5)
            p_pn = tf_pn.paragraphs[0]
            p_pn.text = f"{slide_num}/{total}"
            _set_font(p_pn.runs[0] if p_pn.runs else p_pn.add_run(), 14, False, C_GRAY_TEXT)
            p_pn.alignment = PP_ALIGN.RIGHT

        # Body text
        lines = body.strip().split("\n")
        tf2 = _add_textbox(slide, 1.0, 1.6, 11.333, 5.5)
        for i, line in enumerate(lines):
            clean = line.lstrip("•- ").strip()
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            # Bullet point
            p.text = f"  {clean}"
            _set_font(p.runs[0] if p.runs else p.add_run(), 26, False, C_DARK_TEXT)
            p.space_after = Pt(14)
            # Add bullet marker
            p.text = f"  ●  {clean}"
            # Color the bullet
            if len(p.runs) >= 1:
                pass  # simple bullet styling

        # Bottom line
        slide.shapes.add_shape(
            1, Inches(0), Inches(7.44), SLIDE_W, Inches(0.04)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = C_ACCENT
        slide.shapes[-1].line.fill.background()

    def add_section(self, title: str):
        slide = _blank_slide(self.prs)
        _set_bg(slide, C_DARK_BG)

        # Large centered title
        tf = _add_textbox(slide, 1.5, 2.8, 10.333, 1.8)
        p = tf.paragraphs[0]
        p.text = title
        _set_font(p.runs[0] if p.runs else p.add_run(), 48, True, C_WHITE)
        p.alignment = PP_ALIGN.CENTER

        # Accent line
        slide.shapes.add_shape(
            1, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.04)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = C_ACCENT
        slide.shapes[-1].line.fill.background()

    def add_summary(self, title: str, body: str):
        slide = _blank_slide(self.prs)
        _set_bg(slide, C_DARK_BG)

        tf = _add_textbox(slide, 1.0, 0.8, 11.333, 1.0)
        p = tf.paragraphs[0]
        p.text = title
        _set_font(p.runs[0] if p.runs else p.add_run(), 44, True, C_WHITE)

        # Underline
        slide.shapes.add_shape(
            1, Inches(1.0), Inches(1.7), Inches(3.0), Inches(0.04)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = C_ACCENT
        slide.shapes[-1].line.fill.background()

        lines = body.strip().split("\n")
        tf2 = _add_textbox(slide, 1.5, 2.4, 10.333, 4.5)
        for i, line in enumerate(lines):
            clean = line.lstrip("•- ").strip()
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"  ✓  {clean}"
            _set_font(p.runs[0] if p.runs else p.add_run(), 28, False, C_WHITE)
            p.space_after = Pt(16)

        # Bottom accent
        slide.shapes.add_shape(
            1, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = C_ACCENT
        slide.shapes[-1].line.fill.background()

    def build(self, slides_data: list[dict]) -> Presentation:
        total = len(slides_data)
        for i, s in enumerate(slides_data):
            slide_type = s.get("type", "content")
            if slide_type == "cover":
                self.add_cover(s.get("title", ""), s.get("subtitle", ""))
            elif slide_type == "toc":
                self.add_toc(s.get("title", "目录"), s.get("items", []))
            elif slide_type == "section":
                self.add_section(s.get("title", ""))
            elif slide_type == "summary":
                self.add_summary(s.get("title", "总结"), s.get("body", ""))
            else:  # content
                self.add_content(
                    s.get("title", ""), s.get("body", ""),
                    slide_num=i + 1, total=total,
                )
        return self.prs


# ── External Template Engine ──────────────────────────────────────

TOC_KEYWORDS = ["目录", "目 录", "Contents", "CONTENTS", "目錄"]
LOGO_SKIP_KEYWORDS = ["logo", "LOGO", "http", "www", "@", "//"]


class ExternalTemplate:
    """Fill an external .pptx template.

    Two modes:
    1. Placeholder mode — template uses {{title}} {{body}} etc. in text shapes
    2. Heuristic mode — template has pre-designed slides, engine auto-detects
       types by content/position and replaces text by font-size analysis.
    """

    def __init__(self, template_path: Path, slide_map: Optional[dict] = None):
        self.template_path = template_path
        self.prs = Presentation(str(template_path))
        self.type_map: dict[str, int] = {}       # first index per type
        self.type_pool: dict[str, list[int]] = {}  # all indices per type
        self.has_placeholders = False
        self._shape_cache: dict[int, list[dict]] = {}  # font analysis per src slide

        if slide_map:
            self.type_map = slide_map
            for stype, idx in slide_map.items():
                self.type_pool.setdefault(stype, []).append(idx)
            self.has_placeholders = False
        else:
            self._analyze()

        self._cache_all_slides()

    # ── Analysis ───────────────────────────────────────────────────

    def _analyze(self):
        """Try placeholder detection first, fall back to heuristics."""
        for idx, slide in enumerate(self.prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            full = "\n".join(texts)
            if "{{subtitle}}" in full:
                self.type_map["cover"] = idx
                self.type_pool.setdefault("cover", []).append(idx)
            elif "{{items}}" in full:
                self.type_map["toc"] = idx
                self.type_pool.setdefault("toc", []).append(idx)
            elif "{{body}}" in full:
                self.type_map.setdefault("content", idx)
                self.type_pool.setdefault("content", []).append(idx)
            if "{{section_title}}" in full:
                self.type_map["section"] = idx
                self.type_pool.setdefault("section", []).append(idx)

        if self.type_map:
            self.has_placeholders = True
            logger.info(f"Placeholder mode — types: {list(self.type_map.keys())}")
        else:
            self._analyze_heuristic()
            logger.info(
                f"Heuristic mode — types: {list(self.type_map.keys())} "
                f"pools: { {k: len(v) for k, v in self.type_pool.items()} }"
            )

    def _analyze_heuristic(self):
        """Auto-detect slide types by content keywords, section numbers, position."""
        total = len(self.prs.slides)

        for idx, slide in enumerate(self.prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            full = "\n".join(texts)

            if idx == 0:
                self.type_map["cover"] = idx
                self.type_pool.setdefault("cover", []).append(idx)
                continue

            if idx == total - 1:
                self.type_map["summary"] = idx
                self.type_pool.setdefault("summary", []).append(idx)
                continue

            if any(kw in full for kw in TOC_KEYWORDS):
                self.type_map["toc"] = idx
                self.type_pool.setdefault("toc", []).append(idx)
                continue

            if self._looks_like_section(slide, full):
                self.type_pool.setdefault("section", []).append(idx)
                if "section" not in self.type_map:
                    self.type_map["section"] = idx
                continue

            # Everything else → content pool
            self.type_pool.setdefault("content", []).append(idx)
            if "content" not in self.type_map:
                self.type_map["content"] = idx

        # Remove slides with no usable text shapes from content pool
        for stype in list(self.type_pool.keys()):
            filtered = [
                idx for idx in self.type_pool[stype]
                if self._count_text_shapes(idx) >= (2 if stype == "content" else 1)
            ]
            if filtered:
                self.type_pool[stype] = filtered
                if self.type_map.get(stype) not in filtered:
                    self.type_map[stype] = filtered[0]
            else:
                del self.type_pool[stype]

        # Fallbacks for missing types
        if "toc" not in self.type_map and total >= 2:
            self.type_map["toc"] = 1
            self.type_pool.setdefault("toc", []).append(1)
        if "content" not in self.type_map:
            for idx in range(total):
                if self._count_text_shapes(idx) >= 2:
                    self.type_map["content"] = idx
                    self.type_pool.setdefault("content", []).append(idx)
                    break
        if "summary" not in self.type_map:
            self.type_map["summary"] = 0
            self.type_pool.setdefault("summary", []).append(0)

    def _count_text_shapes(self, idx: int) -> int:
        """Count text shapes on a slide that are not logo/URL elements."""
        slide = self.prs.slides[idx]
        count = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if any(kw in text for kw in LOGO_SKIP_KEYWORDS):
                continue
            count += 1
        return count

    def _looks_like_section(self, slide, full_text: str) -> bool:
        """Heuristic: section divider usually has a section number + page indicator."""
        has_num = bool(re.search(r'\b0[1-9]\b|\b1[0-9]\b', full_text))
        has_page = bool(re.search(r'\d{1,2}\s*/\s*\d{1,2}', full_text))
        text_count = sum(
            1 for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        )
        return has_num and has_page and text_count <= 7

    # ── Shape cache ─────────────────────────────────────────────────

    def _cache_all_slides(self):
        """Pre-analyze text shapes for each template slide (font sizes, skip list)."""
        for idx, slide in enumerate(self.prs.slides):
            infos = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if any(kw in text for kw in LOGO_SKIP_KEYWORDS):
                    continue
                max_font = 0.0
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            max_font = max(max_font, run.font.size)
                infos.append({"max_font": max_font, "text_len": len(text)})
            infos.sort(key=lambda x: x["max_font"], reverse=True)
            self._shape_cache[idx] = infos

    # ── Build ───────────────────────────────────────────────────────

    def build(self, slides_data: list[dict]) -> Presentation:
        """Clone template slides and fill with content."""
        # Clone template via BytesIO — preserves layouts/images without duplicates
        buf = io.BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        out = Presentation(buf)

        # Remove all existing slides, keep layouts
        while out.slides:
            rId = out.slides._sldIdLst[0].get(qn("r:id"))
            out.part.drop_rel(rId)
            out.slides._sldIdLst.remove(out.slides._sldIdLst[0])

        total = len(slides_data)
        type_usage: dict[str, int] = {}

        # Build layout name → index map for output presentation
        out_layout_map = {ly.name: idx for idx, ly in enumerate(out.slide_layouts)}

        for i, s in enumerate(slides_data):
            slide_type = s.get("type", "content")
            pool = self.type_pool.get(slide_type, self.type_pool.get("content", []))
            if not pool:
                logger.warning(f"No template slide for type '{slide_type}', skipping")
                continue

            usage_idx = type_usage.get(slide_type, 0) % len(pool)
            src_idx = pool[usage_idx]
            type_usage[slide_type] = usage_idx + 1

            src_slide = self.prs.slides[src_idx]
            # Use output's own layout to avoid duplicates
            layout_name = src_slide.slide_layout.name
            out_layout_idx = out_layout_map.get(layout_name, 0)
            new_slide = out.slides.add_slide(out.slide_layouts[out_layout_idx])

            for shape in list(new_slide.shapes):
                sp = shape._element
                sp.getparent().remove(sp)
            for shape in src_slide.shapes:
                new_slide.shapes._spTree.append(copy.deepcopy(shape._element))

            data = {
                "title": s.get("title", ""),
                "subtitle": s.get("subtitle", ""),
                "body": s.get("body", ""),
                "items": "\n".join(s.get("items", [])),
                "slide_num": str(i + 1),
                "total_slides": str(total),
            }

            if self.has_placeholders:
                self._fill_placeholders(new_slide, data)
            else:
                self._fill_by_structure(new_slide, src_idx, data, slide_type)

        return out

    # ── Text replacement ────────────────────────────────────────────

    def _fill_placeholders(self, slide, data: dict):
        """Replace {{placeholder}} text in all shapes on the slide."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if not PLACEHOLDER_RE.search(tf.text):
                continue
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = PLACEHOLDER_RE.sub(
                        lambda m: data.get(m.group(1), m.group(0)), run.text
                    )

    def _fill_by_structure(self, slide, src_idx: int, data: dict, slide_type: str):
        """Replace text based on font-size + text-length analysis (no placeholders).

        Strategy:
        - Title = shape with largest font where text_len > 3
        - Subtitle/Body/Items = shape with longest text (not the title)
        - Everything else with text_len > 15 gets cleared.
        """
        entries = []  # (shape, max_font, text_len)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if any(kw in text for kw in LOGO_SKIP_KEYWORDS):
                continue
            max_font = 0.0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_font = max(max_font, run.font.size)
            entries.append((shape, max_font, len(text)))

        if not entries:
            return

        # Title: largest font among shapes with reasonable title length (3-80 chars)
        title_candidates = [(i, e) for i, e in enumerate(entries) if 3 < e[2] <= 80]
        if not title_candidates:
            # No clear title shape → fall back to len > 3
            title_candidates = [(i, e) for i, e in enumerate(entries) if e[2] > 3]
        has_title_shape = bool(title_candidates)
        if not title_candidates:
            title_candidates = list(enumerate(entries))
        title_idx = max(title_candidates, key=lambda x: x[1][1])[0]

        # Content: longest text (not title), only if different from title
        non_title = [(i, e) for i, e in enumerate(entries) if i != title_idx]
        content_idx = max(non_title, key=lambda x: x[1][2])[0] if non_title else title_idx

        # Determine which text to place in the content area
        if slide_type in ("cover", "section", "summary"):
            content_text = data.get("subtitle", "") or data.get("body", "")
        elif slide_type == "toc":
            content_text = data.get("items", "")
        else:
            content_text = data.get("body", "")

        # For content slides without a real title shape, prepend title to body
        body_with_title = content_text
        if slide_type == "content" and not has_title_shape and data.get("title") and content_text:
            body_with_title = data["title"] + "\n" + content_text

        # Only protect content_idx if we actually fill it with something
        protected = {title_idx}
        if content_text and content_idx != title_idx:
            protected.add(content_idx)
        if slide_type == "content" and not has_title_shape:
            # Don't protect title_idx if it's not a real title shape
            protected.discard(title_idx)

        # Fill title
        if has_title_shape or slide_type not in ("content",):
            self._replace_shape_text(entries[title_idx][0], data.get("title", ""))

        # Fill content area
        if content_idx != title_idx and body_with_title:
            self._replace_shape_text(entries[content_idx][0], body_with_title)

        # Clear unprotected shapes (remove stale template filler)
        for j, (shape, _, tlen) in enumerate(entries):
            if j in protected:
                continue
            old_text = shape.text_frame.text.strip()
            if re.match(r'^\d{1,2}$', old_text):
                self._replace_shape_text(shape, data.get("slide_num", old_text))
            elif re.match(r'^\d{1,3}\s*/\s*\d{1,3}$', old_text):
                self._replace_shape_text(shape, data.get("slide_num", ""))
            elif tlen > 4:
                self._replace_shape_text(shape, "")

    def _replace_shape_text(self, shape, new_text: str):
        """Replace all text in a shape, keeping formatting of the first run."""
        tf = shape.text_frame
        if not tf.paragraphs:
            return
        # Remove extra paragraphs
        for para in list(tf.paragraphs)[1:]:
            para._p.getparent().remove(para._p)
        first = tf.paragraphs[0]
        if first.runs:
            first.runs[0].text = new_text
            for run in first.runs[1:]:
                run.text = ""
        else:
            first.add_run().text = new_text


# ── Public API ────────────────────────────────────────────────────

def generate_pptx(slides_data: list[dict], output_path: Path,
                  template_path: Optional[str] = None,
                  slide_map: Optional[dict] = None) -> Path:
    """Main entry: choose template engine and build PPTX.

    Args:
        slides_data: List of {"type": "cover"|"toc"|"content"|"section"|"summary", ...}
        output_path: Destination .pptx path.
        template_path: Optional external .pptx template.
        slide_map: Optional manual mapping e.g. {"cover": 0, "toc": 1, "content": 3}.

    Returns:
        The output_path on success.
    """
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if template_path:
        tp = Path(template_path)
        if tp.exists():
            logger.info(f"Using external template: {tp}")
            engine = ExternalTemplate(tp, slide_map=slide_map)
            prs = engine.build(slides_data)
        else:
            logger.warning(f"Template not found: {tp}, using built-in")
            prs = BuiltinTemplate().build(slides_data)
    else:
        logger.info("Using built-in template")
        prs = BuiltinTemplate().build(slides_data)

    prs.save(str(output_path))
    logger.info(f"PPTX saved: {output_path} ({len(slides_data)} slides)")
    return output_path
